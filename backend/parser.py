"""
parser.py
Extracts plain text from uploaded files so it can be sent to the AI.

Supported: .pdf, .docx, .pptx, .txt, .md
Each extractor returns a single cleaned string. Add OCR (pytesseract) later
if you want image support.
"""

import os
import pdfplumber
from docx import Document
from pptx import Presentation


class UnsupportedFileType(Exception):
    pass


def extract_pdf(filepath: str) -> str:
    text_chunks = []
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_chunks.append(page_text)
    return "\n".join(text_chunks)


def extract_docx(filepath: str) -> str:
    doc = Document(filepath)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_pptx(filepath: str) -> str:
    prs = Presentation(filepath)
    text_chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_chunks.append(shape.text)
    return "\n".join(text_chunks)


def extract_txt(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".txt": extract_txt,
    ".md": extract_txt,
}


def extract_text(filepath: str) -> str:
    """
    Detects file type by extension and routes to the right extractor.
    Cleans up excess whitespace before returning.
    """
    ext = os.path.splitext(filepath)[1].lower()
    extractor = EXTRACTORS.get(ext)

    if extractor is None:
        raise UnsupportedFileType(
            f"'{ext}' files are not supported. Allowed: {list(EXTRACTORS.keys())}"
        )

    raw_text = extractor(filepath)
    cleaned = "\n".join(line.strip() for line in raw_text.splitlines() if line.strip())

    if not cleaned:
        raise ValueError("No readable text could be extracted from this file.")

    return cleaned


def allowed_file(filename: str) -> bool:
    ext = os.path.splitext(filename)[1].lower()
    return ext in EXTRACTORS
